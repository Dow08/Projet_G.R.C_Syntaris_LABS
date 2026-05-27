#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
generate_pptx.py
----------------
Génère une présentation PowerPoint (.pptx) de niveau professionnel
pour Dorian Poncelet sur l'audit GRC de Syntaris Group.
Design épuré, asymétrique, conforme aux meilleures pratiques GRC et de design.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Couleurs de la charte graphique : Midnight Executive & Teal Trust
COLOR_DARK_NAVY = RGBColor(10, 17, 40)      # Arrière-plan sombre (60%)
COLOR_LIGHT_GRAY = RGBColor(245, 247, 250)   # Arrière-plan clair (60%)
COLOR_CARD_WHITE = RGBColor(255, 255, 255)   # Cartes de contenu
COLOR_TEXT_DARK = RGBColor(15, 23, 42)       # Texte sur fond clair
COLOR_TEXT_LIGHT = RGBColor(241, 245, 249)   # Texte sur fond sombre
COLOR_ACCENT_TEAL = RGBColor(2, 128, 144)    # Accent principal (Confiance)
COLOR_ACCENT_CORAL = RGBColor(249, 97, 103)   # Accent secondaire (Alerte / Écart)
COLOR_MUTED_GRAY = RGBColor(100, 116, 139)   # Texte secondaire / captions

def appliquer_fond_sombre(slide):
    """Applique un arrière-plan bleu nuit uni à la diapositive."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_DARK_NAVY

def appliquer_fond_clair(slide):
    """Applique un arrière-plan gris clair uni à la diapositive."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_LIGHT_GRAY

def ajouter_titre_diapo(slide, texte):
    """Ajoute un titre standard aligné à gauche sur une diapositive claire."""
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.83), Inches(0.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = texte
    p.font.name = 'Trebuchet MS'
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLOR_DARK_NAVY

def creer_carte_blanche(slide, x, y, w, h):
    """Crée un rectangle blanc servant de conteneur visuel (carte)."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_CARD_WHITE
    shape.line.color.rgb = RGBColor(226, 232, 240)  # Ligne grise très claire
    shape.line.width = Pt(1)
    return shape

# =====================================================================
# INITIALISATION DE LA PRESENTATION
# =====================================================================
prs = Presentation()
prs.slide_width = Inches(13.333)   # Format 16:9 widescreen
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6] # Layout vide pour design sur mesure

# ---------------------------------------------------------------------
# SLIDE 1 : TITRE (FOND SOMBRE)
# ---------------------------------------------------------------------
slide1 = prs.slides.add_slide(blank_layout)
appliquer_fond_sombre(slide1)

# Zone de titre principal
title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.33), Inches(2.5))
tf = title_box.text_frame
tf.word_wrap = True
tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

p1 = tf.paragraphs[0]
p1.text = "AUDIT GRC & MATURITÉ SÉCURITÉ"
p1.font.name = 'Trebuchet MS'
p1.font.size = Pt(44)
p1.font.bold = True
p1.font.color.rgb = COLOR_TEXT_LIGHT

p2 = tf.add_paragraph()
p2.text = "Plan Stratégique Conjoint NIST CSF 2.0 & RGPD"
p2.font.name = 'Georgia'
p2.font.size = Pt(24)
p2.font.color.rgb = COLOR_ACCENT_TEAL
p2.space_before = Pt(15)

# Zone de sous-titre / auteur
author_box = slide1.shapes.add_textbox(Inches(1.0), Inches(5.2), Inches(11.33), Inches(1.2))
tf_auth = author_box.text_frame
tf_auth.word_wrap = True
tf_auth.margin_left = tf_auth.margin_right = tf_auth.margin_top = tf_auth.margin_bottom = 0

pa1 = tf_auth.paragraphs[0]
pa1.text = "Présenté par : Dorian Poncelet"
pa1.font.name = 'Calibri'
pa1.font.size = Pt(16)
pa1.font.bold = True
pa1.font.color.rgb = COLOR_TEXT_LIGHT

pa2 = tf_auth.add_paragraph()
pa2.text = "Conseiller externe en Cybersécurité | Dossier Syntaris Group"
pa2.font.name = 'Calibri'
pa2.font.size = Pt(14)
pa2.font.color.rgb = COLOR_MUTED_GRAY
pa2.space_before = Pt(5)


# ---------------------------------------------------------------------
# SLIDE 2 : CONTEXTE DE L'AUDIT (FOND CLAIR)
# ---------------------------------------------------------------------
slide2 = prs.slides.add_slide(blank_layout)
appliquer_fond_clair(slide2)
ajouter_titre_diapo(slide2, "Contexte de l'Audit : Syntaris Group")

# Carte de gauche : Profil de l'entreprise
creer_carte_blanche(slide2, 0.75, 1.5, 5.6, 5.2)
p_box1 = slide2.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(5.1), Inches(4.8))
tf1 = p_box1.text_frame
tf1.word_wrap = True
tf1.margin_left = tf1.margin_right = tf1.margin_top = tf1.margin_bottom = 0

hp1 = tf1.paragraphs[0]
hp1.text = "PROFIL DE L'ORGANISATION"
hp1.font.name = 'Trebuchet MS'
hp1.font.size = Pt(18)
hp1.font.bold = True
hp1.font.color.rgb = COLOR_ACCENT_TEAL

bp1 = tf1.add_paragraph()
bp1.text = "• Activité : Fintech de paiement en ligne et de vérification d'identité numérique (KYC).\n\n• Envergure : Croissance rapide, opérations transatlantiques (Europe & États-Unis).\n\n• Architecture cible : Hybride Cloud-Local. Plateforme transactionnelle sur AWS/Azure ; postes administratifs et corporatifs sur site."
bp1.font.name = 'Calibri'
bp1.font.size = Pt(15)
bp1.font.color.rgb = COLOR_TEXT_DARK
bp1.space_before = Pt(15)

# Carte de droite : Enjeux et Facteurs Réglementaires
creer_carte_blanche(slide2, 6.98, 1.5, 5.6, 5.2)
p_box2 = slide2.shapes.add_textbox(Inches(7.23), Inches(1.7), Inches(5.1), Inches(4.8))
tf2 = p_box2.text_frame
tf2.word_wrap = True
tf2.margin_left = tf2.margin_right = tf2.margin_top = tf2.margin_bottom = 0

hp2 = tf2.paragraphs[0]
hp2.text = "ENJEUX STRATÉGIQUES & EXIGENCES"
hp2.font.name = 'Trebuchet MS'
hp2.font.size = Pt(18)
hp2.font.bold = True
hp2.font.color.rgb = COLOR_ACCENT_TEAL

bp2 = tf2.add_paragraph()
bp2.text = "• Régulations : Soumis de plein droit au RGPD européen et aux exigences financières PCI-DSS pour le paiement.\n\n• Exigence Marché : Les banques et marchands partenaires exigent des preuves solides d'une gouvernance cyber mature (NIST CSF 2.0) avant de contractualiser.\n\n• Objectif : Structurer une trajectoire réaliste de remédiation."
bp2.font.name = 'Calibri'
bp2.font.size = Pt(15)
bp2.font.color.rgb = COLOR_TEXT_DARK
bp2.space_before = Pt(15)


# ---------------------------------------------------------------------
# SLIDE 3 : CARTOGRAPHIE DES ACTIFS (FOND CLAIR)
# ---------------------------------------------------------------------
slide3 = prs.slides.add_slide(blank_layout)
appliquer_fond_clair(slide3)
ajouter_titre_diapo(slide3, "Cartographie des Actifs & Données Critiques")

# 3 Colonnes asymétriques pour les types de données
col_w = 3.6
col_h = 5.2
col_gap = 0.51
start_x = 0.75

donnees_types = [
    {
        "titre": "DONNÉES FINANCIÈRES",
        "ss_titre": "Conformité PCI-DSS",
        "desc": "• Numéros de cartes (PAN)\n• Historique des transactions\n• Coordonnées bancaires commerçants\n\nExige un chiffrement de bout en bout et une isolation stricte des bases de données transactionnelles."
    },
    {
        "titre": "DONNÉES BIOMÉTRIQUES",
        "ss_titre": "RGPD Article 9 (Sensible)",
        "desc": "• Scans de passeports/ID\n• Gabarits de reconnaissance faciale pour vérification KYC\n\nLa collecte de ces données sensibles rend strictement obligatoire la désignation d'un DPO et une Analyse d'Impact (AIPD)."
    },
    {
        "titre": "PROPRIÉTÉ & SYSTÈMES",
        "ss_titre": "Actifs Intellectuels",
        "desc": "• Code source des APIs de paiement\n• Clés API et credentials Cloud\n• Bases d'identités Active Directory\n\nNécessite des pipelines de déploiement sécurisés (CI/CD) et une authentification forte pour les développeurs."
    }
]

for i, dtype in enumerate(donnees_types):
    cx = start_x + i * (col_w + col_gap)
    creer_carte_blanche(slide3, cx, 1.5, col_w, col_h)
    
    box = slide3.shapes.add_textbox(Inches(cx + 0.2), Inches(1.7), Inches(col_w - 0.4), Inches(col_h - 0.4))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    
    hp = tf.paragraphs[0]
    hp.text = dtype["titre"]
    hp.font.name = 'Trebuchet MS'
    hp.font.size = Pt(18)
    hp.font.bold = True
    hp.font.color.rgb = COLOR_ACCENT_TEAL
    
    sub = tf.add_paragraph()
    sub.text = dtype["ss_titre"]
    sub.font.name = 'Georgia'
    sub.font.size = Pt(13)
    sub.font.italic = True
    sub.font.color.rgb = COLOR_MUTED_GRAY
    sub.space_before = Pt(5)
    
    bp = tf.add_paragraph()
    bp.text = dtype["desc"]
    bp.font.name = 'Calibri'
    bp.font.size = Pt(14)
    bp.font.color.rgb = COLOR_TEXT_DARK
    bp.space_before = Pt(20)


# ---------------------------------------------------------------------
# SLIDE 4 : LES TROIS RISQUES MAJEURS (FOND CLAIR)
# ---------------------------------------------------------------------
slide4 = prs.slides.add_slide(blank_layout)
appliquer_fond_clair(slide4)
ajouter_titre_diapo(slide4, "Top 3 des Risques d'Entreprise")

risques = [
    {
        "num": "01",
        "titre": "FUITE DE LA BASE KYC ET BIOMÉTRIQUE",
        "contexte": "Le stockage indéfini de passeports et empreintes faciales sans isolation forte expose Syntaris à un piratage massif.",
        "impact": "Impact financier direct : Amendes CNIL (jusqu'à 4% du CA mondial - Art. 83), poursuites judiciaires collectives, perte totale de confiance commerciale."
    },
    {
        "num": "02",
        "titre": "INDISPONIBILITÉ DU SERVICE DE PAIEMENT",
        "contexte": "Une attaque DDoS d'envergure ou un rançongiciel (ransomware) paralyse le moteur de paiement hébergé sur le Cloud public.",
        "impact": "Impact opérationnel immédiat : Rupture des engagements de service (SLA) avec les commerçants, pertes de transactions sèches par minute."
    },
    {
        "num": "03",
        "titre": "PIRATAGE DE L'API DE PAIEMENT (SUPPLY CHAIN)",
        "contexte": "La compromission du poste d'un développeur permet l'injection de code espion (type Magecart) au sein même du code des transactions.",
        "impact": "Impact réputationnel fatal : Vol des numéros de carte de paiement des clients finaux en temps réel lors du traitement des paiements."
    }
]

for i, r in enumerate(risques):
    ry = 1.5 + i * 1.8
    # Carte blanche pour chaque risque
    creer_carte_blanche(slide4, 0.75, ry, 11.83, 1.5)
    
    # Numéro d'alerte en corail
    num_box = slide4.shapes.add_textbox(Inches(0.95), Inches(ry + 0.15), Inches(1.0), Inches(1.2))
    tf_num = num_box.text_frame
    p_num = tf_num.paragraphs[0]
    p_num.text = r["num"]
    p_num.font.name = 'Trebuchet MS'
    p_num.font.size = Pt(40)
    p_num.font.bold = True
    p_num.font.color.rgb = COLOR_ACCENT_CORAL
    
    # Texte du risque
    text_box = slide4.shapes.add_textbox(Inches(2.0), Inches(ry + 0.15), Inches(10.2), Inches(1.2))
    tf_text = text_box.text_frame
    tf_text.word_wrap = True
    tf_text.margin_left = tf_text.margin_right = tf_text.margin_top = tf_text.margin_bottom = 0
    
    p_title = tf_text.paragraphs[0]
    p_title.text = r["titre"]
    p_title.font.name = 'Trebuchet MS'
    p_title.font.size = Pt(15)
    p_title.font.bold = True
    p_title.font.color.rgb = COLOR_DARK_NAVY
    
    p_desc = tf_text.add_paragraph()
    p_desc.text = f"{r['contexte']}\n{r['impact']}"
    p_desc.font.name = 'Calibri'
    p_desc.font.size = Pt(13)
    p_desc.font.color.rgb = COLOR_TEXT_DARK
    p_desc.space_before = Pt(5)


# ---------------------------------------------------------------------
# SLIDE 5 : PROFIL DE MATURITE NIST CSF 2.0 (FOND CLAIR)
# ---------------------------------------------------------------------
slide5 = prs.slides.add_slide(blank_layout)
appliquer_fond_clair(slide5)
ajouter_titre_diapo(slide5, "Maturité NIST CSF 2.0 : Profil Actuel vs Cible")

# Tableau récapitulatif
rows = 8
cols = 5
table_shape = slide5.shapes.add_table(rows, cols, Inches(0.75), Inches(1.5), Inches(11.83), Inches(5.2))
table = table_shape.table

# Largeur des colonnes
table.columns[0].width = Inches(2.33)  # Fonction
table.columns[1].width = Inches(1.5)   # Score Actuel
table.columns[2].width = Inches(1.5)   # Score Cible
table.columns[3].width = Inches(1.5)   # Écart
table.columns[4].width = Inches(5.0)   # Constat Majeur d'Audit

headers = ["Fonction NIST CSF 2.0", "Score Actuel", "Score Cible", "Écart (Gap)", "Observation majeure de l'Audit"]
for j, h in enumerate(headers):
    cell = table.cell(0, j)
    cell.text = h
    cell.fill.solid()
    cell.fill.fore_color.rgb = COLOR_DARK_NAVY
    for p in cell.text_frame.paragraphs:
        p.font.name = 'Trebuchet MS'
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = COLOR_TEXT_LIGHT

eval_data = [
    ["Gouverner (Govern - GV)", "1.50 / 4.0", "3.67 / 4.0", "+2.17", "Pas de Politique de Sécurité (PSSI) globale formulée par la direction."],
    ["Identifier (Identify - ID)", "1.50 / 4.0", "3.75 / 4.0", "+2.25", "Aucune cartographie documentée des flux de données de paiement et KYC."],
    ["Protéger (Protect - PR)", "1.83 / 4.0", "4.00 / 4.0", "+2.17", "HTTPS présent mais MFA absent sur les consoles et postes des développeurs."],
    ["Détecter (Detect - DE)", "1.50 / 4.0", "4.00 / 4.0", "+2.50", "Logs cloud existants mais pas de SIEM ni d'équipe SOC pour lever les doutes."],
    ["Répondre (Respond - RS)", "1.00 / 4.0", "4.00 / 4.0", "+3.00", "ZÉRO PLAYBOOK : Aucun plan de réponse en cas d'intrusion ou fuite (Art 33)."],
    ["Récupérer (Recover - RC)", "1.00 / 4.0", "4.00 / 4.0", "+3.00", "ZÉRO PCA/PRA : Aucun test régulier de restauration des sauvegardes cloud."],
]

for i, row_data in enumerate(eval_data):
    for j, val in enumerate(row_data):
        cell = table.cell(i + 1, j)
        cell.text = val
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_CARD_WHITE
        for p in cell.text_frame.paragraphs:
            p.font.name = 'Calibri'
            p.font.size = Pt(12)
            p.font.color.rgb = COLOR_TEXT_DARK
            # Formater les écarts critiques en gras corail
            if j == 3 and "+" in val:
                p.font.bold = True
                p.font.color.rgb = COLOR_ACCENT_CORAL
            # Mettre en gras le nom des fonctions
            if j == 0:
                p.font.bold = True

# Ligne de synthèse globale tout en bas
cell_tot_name = table.cell(7, 0)
cell_tot_name.text = "SYNTHÈSE GLOBALE"
cell_tot_act = table.cell(7, 1)
cell_tot_act.text = "1.50 / 4.0"
cell_tot_cib = table.cell(7, 2)
cell_tot_cib.text = "3.86 / 4.0"
cell_tot_gap = table.cell(7, 3)
cell_tot_gap.text = "+2.36"
cell_tot_obs = table.cell(7, 4)
cell_tot_obs.text = "Maturité faible (Tier 1/2) exigeant de profonds ajustements de gouvernance."

for j in range(5):
    cell = table.cell(7, j)
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(241, 245, 249) # Fond gris clair de synthèse
    for p in cell.text_frame.paragraphs:
        p.font.name = 'Calibri'
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = COLOR_TEXT_DARK
        if j == 3:
            p.font.color.rgb = COLOR_ACCENT_CORAL


# ---------------------------------------------------------------------
# SLIDE 6 : RÉSULTATS D'AUDIT RGPD (FOND CLAIR)
# ---------------------------------------------------------------------
slide6 = prs.slides.add_slide(blank_layout)
appliquer_fond_clair(slide6)
ajouter_titre_diapo(slide6, "Audit RGPD : Les Écarts Critiques")

gaps_rgpd = [
    {
        "art": "Art 5",
        "nom": "Minimisation & Conservation",
        "constat": "Les passeports scannés pour la validation KYC biométrique sont conservés indéfiniment en base de données.",
        "solution": "Purge automatique de l'image du passeport sous 30 jours ; conservation exclusive du jeton de validation numérique."
    },
    {
        "art": "Art 25",
        "nom": "Privacy by Design & Default",
        "constat": "Les développeurs utilisent de vraies données clients dans les environnements de test locaux.",
        "solution": "Masquage de données et génération de jeux de tests synthétiques et fictifs par défaut pour les équipes dev."
    },
    {
        "art": "Art 32",
        "nom": "Sécurité des Traitements",
        "constat": "Le MFA n'est pas obligatoire sur les dépôts de code (GitHub) ; bases de données cloud non isolées.",
        "solution": "Obligation du MFA, intégration de scans de vulnérabilités applicatifs (SAST) dans la pipeline CI/CD."
    },
    {
        "art": "Art 33",
        "nom": "Notification sous 72 Heures",
        "constat": "Aucun DPO n'est nommé et aucune procédure de déclaration de fuite de données n'existe pour la CNIL.",
        "solution": "Nomination d'un DPO et écriture immédiate d'un protocole d'escalade d'incident d'intrusion."
    }
]

card_w = 2.7
card_h = 5.2
card_gap = 0.34
start_x = 0.75

for i, g in enumerate(gaps_rgpd):
    cx = start_x + i * (card_w + card_gap)
    creer_carte_blanche(slide6, cx, 1.5, card_w, card_h)
    
    # Insérer le contenu
    box = slide6.shapes.add_textbox(Inches(cx + 0.15), Inches(1.7), Inches(card_w - 0.3), Inches(card_h - 0.4))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    
    hp = tf.paragraphs[0]
    hp.text = g["art"]
    hp.font.name = 'Trebuchet MS'
    hp.font.size = Pt(24)
    hp.font.bold = True
    hp.font.color.rgb = COLOR_ACCENT_CORAL
    
    sub = tf.add_paragraph()
    sub.text = g["nom"]
    sub.font.name = 'Trebuchet MS'
    sub.font.size = Pt(14)
    sub.font.bold = True
    sub.font.color.rgb = COLOR_DARK_NAVY
    sub.space_before = Pt(5)
    
    bp1 = tf.add_paragraph()
    bp1.text = f"CONSTAT D'AUDIT :\n{g['constat']}"
    bp1.font.name = 'Calibri'
    bp1.font.size = Pt(11)
    bp1.font.color.rgb = COLOR_TEXT_DARK
    bp1.space_before = Pt(15)
    
    bp2 = tf.add_paragraph()
    bp2.text = f"REMÉDIATION CONCRÈTE :\n{g['solution']}"
    bp2.font.name = 'Calibri'
    bp2.font.size = Pt(11)
    bp2.font.bold = True
    bp2.font.color.rgb = COLOR_ACCENT_TEAL
    bp2.space_before = Pt(10)


# ---------------------------------------------------------------------
# SLIDE 7 : FEUILLE DE ROUTE - PRIORITE 1 (FOND CLAIR)
# ---------------------------------------------------------------------
slide7 = prs.slides.add_slide(blank_layout)
appliquer_fond_clair(slide7)
ajouter_titre_diapo(slide7, "Roadmap Priorité 1 : Quick Wins (0 - 30 Jours)")

# Deux blocs horizontaux distincts
creer_carte_blanche(slide7, 0.75, 1.5, 11.83, 2.3)
box_rec1 = slide7.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(11.33), Inches(2.0))
tf_rec1 = box_rec1.text_frame
tf_rec1.word_wrap = True
tf_rec1.margin_left = tf_rec1.margin_right = tf_rec1.margin_top = tf_rec1.margin_bottom = 0

hp_rec1 = tf_rec1.paragraphs[0]
hp_rec1.text = "REC-01 : APPOINTEMENT D'UN DPO & SÉCURISATION DU PLAN D'URGENCE"
hp_rec1.font.name = 'Trebuchet MS'
hp_rec1.font.size = Pt(16)
hp_rec1.font.bold = True
hp_rec1.font.color.rgb = COLOR_ACCENT_CORAL

bp_rec1 = tf_rec1.add_paragraph()
bp_rec1.text = "• Obligation légale : Rendre Syntaris conforme de toute urgence aux obligations de DPO (Art 37 RGPD) face aux volumes biométriques.\n• Protocole 72H : Rédiger un document simple (playbook) listant les étapes d'escalade d'une intrusion et pré-remplissant le document de déclaration officielle CNIL en cas de fuite de données.\n• Responsable : Service Juridique / CEO."
bp_rec1.font.name = 'Calibri'
bp_rec1.font.size = Pt(14)
bp_rec1.font.color.rgb = COLOR_TEXT_DARK
bp_rec1.space_before = Pt(10)

creer_carte_blanche(slide7, 0.75, 4.2, 11.83, 2.3)
box_rec2 = slide7.shapes.add_textbox(Inches(1.0), Inches(4.4), Inches(11.33), Inches(2.0))
tf_rec2 = box_rec2.text_frame
tf_rec2.word_wrap = True
tf_rec2.margin_left = tf_rec2.margin_right = tf_rec2.margin_top = tf_rec2.margin_bottom = 0

hp_rec2 = tf_rec2.paragraphs[0]
hp_rec2.text = "REC-02 : MULTI-FACTEUR (MFA) SYSTEMATIQUE ET EXCLUSIF"
hp_rec2.font.name = 'Trebuchet MS'
hp_rec2.font.size = Pt(16)
hp_rec2.font.bold = True
hp_rec2.font.color.rgb = COLOR_ACCENT_CORAL

bp_rec2 = tf_rec2.add_paragraph()
bp_rec2.text = "• Cloud public : Bloquer l'accès aux consoles de management AWS/Azure de production par authentification forte MFA obligatoire.\n• Code source : Forcer l'utilisation de clés MFA matérielles pour l'ensemble des développeurs de la fintech accédant aux bases GitHub.\n• Responsable : Équipe DevOps / Administrateur des Systèmes."
bp_rec2.font.name = 'Calibri'
bp_rec2.font.size = Pt(14)
bp_rec2.font.color.rgb = COLOR_TEXT_DARK
bp_rec2.space_before = Pt(10)


# ---------------------------------------------------------------------
# SLIDE 8 : FEUILLE DE ROUTE - PRIORITES 2 & 3 (FOND CLAIR)
# ---------------------------------------------------------------------
slide8 = prs.slides.add_slide(blank_layout)
appliquer_fond_clair(slide8)
ajouter_titre_diapo(slide8, "Roadmap Priorités 2 & 3 : Chantiers à Moyen Terme")

# Colonne de gauche : Priorité 2 (30-90 Jours)
creer_carte_blanche(slide8, 0.75, 1.5, 5.6, 5.2)
box_mid1 = slide8.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(5.1), Inches(4.8))
tf_mid1 = box_mid1.text_frame
tf_mid1.word_wrap = True
tf_mid1.margin_left = tf_mid1.margin_right = tf_mid1.margin_top = tf_mid1.margin_bottom = 0

hp_mid1 = tf_mid1.paragraphs[0]
hp_mid1.text = "PRIORITÉ 2 : ARCHITECTURE TECHNIQUE (30-90 Jours)"
hp_mid1.font.name = 'Trebuchet MS'
hp_mid1.font.size = Pt(16)
hp_mid1.font.bold = True
hp_mid1.font.color.rgb = COLOR_ACCENT_TEAL

bp_mid1 = tf_mid1.add_paragraph()
bp_mid1.text = "• REC-03 : Tokenisation PCI-DSS & Données Fictives (Art 25)\n  Intégrer une passerelle de tokenisation pour sortir les données de cartes des bases de données et automatiser le masquage pour les tests dev.\n\n• REC-04 : Chiffrement AES-256 & Sauvegardes Froides\n  Chiffrer au repos l'ensemble des bases d'identité cloud et conserver un jeu de sauvegardes régulier déconnecté et immuable face aux ransomwares."
bp_mid1.font.name = 'Calibri'
bp_mid1.font.size = Pt(14)
bp_mid1.font.color.rgb = COLOR_TEXT_DARK
bp_mid1.space_before = Pt(15)

# Colonne de droite : Priorité 3 (90-180 Jours)
creer_carte_blanche(slide8, 6.98, 1.5, 5.6, 5.2)
box_mid2 = slide8.shapes.add_textbox(Inches(7.23), Inches(1.7), Inches(5.1), Inches(4.8))
tf_mid2 = box_mid2.text_frame
tf_mid2.word_wrap = True
tf_mid2.margin_left = tf_mid2.margin_right = tf_mid2.margin_top = tf_mid2.margin_bottom = 0

hp_mid2 = tf_mid2.paragraphs[0]
hp_mid2.text = "PRIORITÉ 3 : PILOTAGE & CULTURE (90-180 Jours)"
hp_mid2.font.name = 'Trebuchet MS'
hp_mid2.font.size = Pt(16)
hp_mid2.font.bold = True
hp_mid2.font.color.rgb = COLOR_ACCENT_TEAL

bp_mid2 = tf_mid2.add_paragraph()
bp_mid2.text = "• REC-05 : Rédaction de la PSSI & Charte Informatique\n  Documenter l'ensemble des règles de sécurité imposées par la fintech. Imposer la signature électronique de la charte utilisateur aux salariés.\n\n• REC-06 : Surveillance 24/7 par SOC Managé\n  Centraliser l'ensemble des logs de transactions dans un SIEM et externaliser la surveillance de sécurité en temps réel à un SOC managé tiers."
bp_mid2.font.name = 'Calibri'
bp_mid2.font.size = Pt(14)
bp_mid2.font.color.rgb = COLOR_TEXT_DARK
bp_mid2.space_before = Pt(15)


# ---------------------------------------------------------------------
# SLIDE 9 : CONCLUSION (FOND SOMBRE)
# ---------------------------------------------------------------------
slide9 = prs.slides.add_slide(blank_layout)
appliquer_fond_sombre(slide9)

# Message final fort
title_box = slide9.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.33), Inches(3.0))
tf = title_box.text_frame
tf.word_wrap = True
tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

p1 = tf.paragraphs[0]
p1.text = "SÉCURITÉ CYBER ET GRC : UN COMPROMIS MAJEUR"
p1.font.name = 'Trebuchet MS'
p1.font.size = Pt(36)
p1.font.bold = True
p1.font.color.rgb = COLOR_TEXT_LIGHT

p2 = tf.add_paragraph()
p2.text = "• Conformité comme levier d'affaires : Remplir les exigences NIST CSF 2.0 est un passeport obligatoire pour contractualiser avec les institutions bancaires et grands marchands.\n\n• Sécurité active de la donnée : Éviter le sur-stockage biométrique et automatiser la conformité par des scripts techniques protège Syntaris des risques financiers les plus critiques.\n\n• Démarche GRC technique : Démontrer la sécurité par des données structurées et du code apporte une valeur immense et différenciante."
p2.font.name = 'Georgia'
p2.font.size = Pt(16)
p2.font.color.rgb = COLOR_TEXT_LIGHT
p2.space_before = Pt(25)

# Pied de page conclusion
end_box = slide9.shapes.add_textbox(Inches(1.0), Inches(5.8), Inches(11.33), Inches(1.0))
tf_end = end_box.text_frame
p_end = tf_end.paragraphs[0]
p_end.text = "Merci pour votre attention. Questions & Échanges."
p_end.font.name = 'Trebuchet MS'
p_end.font.size = Pt(20)
p_end.font.bold = True
p_end.font.color.rgb = COLOR_ACCENT_TEAL

# Sauvegarde de la présentation
output_pptx_path = "/workspace/reports/presentation_slides.pptx"
prs.save(output_pptx_path)
print(f"[+] Présentation PowerPoint professionnelle générée avec succès dans : {output_pptx_path}")
